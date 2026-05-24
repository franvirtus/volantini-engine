from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
TEMPLATE_PATH = ROOT / "templates" / "combat_light_vertical.pptx"

SLIDE_W = 8.27
SLIDE_H = 11.69

WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(183, 255, 0)
GREEN_DARK = RGBColor(142, 214, 0)
BLACK = RGBColor(5, 5, 5)
GRAY = RGBColor(242, 242, 242)
TEXT_GRAY = RGBColor(68, 68, 68)
LIGHT_GREEN = RGBColor(210, 255, 92)

DISPLAY_FONT = "Bebas Neue"
DISPLAY_FALLBACK = "Anton"
BODY_FONT = "Montserrat"


def inch(value):
    return Inches(value)


def set_no_line(shape):
    shape.line.fill.background()


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def add_text(slide, x, y, w, h, text, size, color=BLACK, font=BODY_FONT,
             bold=False, align=PP_ALIGN.LEFT, name=None, rotation=0):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    if name:
        box.name = name
    box.rotation = rotation
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_placeholder(slide, name, x, y, w, h, label=None, fill=GRAY):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        inch(x),
        inch(y),
        inch(w),
        inch(h),
    )
    shape.name = name
    set_fill(shape, fill)
    shape.line.color.rgb = RGBColor(218, 218, 218)
    shape.line.width = Pt(1)

    return shape


def add_skew_bar(slide, x, y, w, h, color=GREEN, rotation=-11, name=None):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, inch(x), inch(y), inch(w), inch(h))
    if name:
        bar.name = name
    bar.rotation = rotation
    set_fill(bar, color)
    set_no_line(bar)
    return bar


def add_logo_mark_placeholder(slide, name, x, y, size):
    holder = add_placeholder(slide, name, x, y, size, size, "", WHITE)
    holder.line.color.rgb = GREEN
    holder.line.width = Pt(1.5)
    v1 = add_skew_bar(slide, x + size * 0.12, y + size * 0.08, size * 0.2, size * 0.72, GREEN, -22)
    v2 = add_skew_bar(slide, x + size * 0.38, y + size * 0.08, size * 0.16, size * 0.72, GREEN, 22)
    v3 = add_skew_bar(slide, x + size * 0.58, y + size * 0.08, size * 0.14, size * 0.55, GREEN, 22)
    v1.name = f"{name}_editable_mark_1"
    v2.name = f"{name}_editable_mark_2"
    v3.name = f"{name}_editable_mark_3"
    return holder


def add_age_box(slide):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
        inch(0.28),
        inch(3.92),
        inch(3.78),
        inch(0.62),
    )
    box.name = "AGE_BORDER"
    box.rotation = -3
    box.fill.background()
    box.line.color.rgb = GREEN
    box.line.width = Pt(1.5)
    add_text(
        slide,
        0.55,
        4.0,
        3.18,
        0.36,
        "{{AGE}}",
        23,
        BLACK,
        DISPLAY_FALLBACK,
        True,
        PP_ALIGN.CENTER,
        "TEXT_AGE",
        -3,
    )


def add_benefit_icon(slide, index, cx, y):
    green = GREEN_DARK
    name = f"BENEFIT_{index}_ICON"
    if index == 1:
        outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - 0.22), inch(y), inch(0.44), inch(0.44))
        outer.name = name
        outer.fill.background()
        outer.line.color.rgb = green
        outer.line.width = Pt(2)
        inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - 0.12), inch(y + 0.10), inch(0.24), inch(0.24))
        inner.fill.background()
        inner.line.color.rgb = green
        inner.line.width = Pt(1.5)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - 0.03), inch(y + 0.19), inch(0.06), inch(0.06))
        set_fill(dot, green)
        dot.line.color.rgb = green
    elif index == 2:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, inch(cx - 0.22), inch(y - 0.01), inch(0.44), inch(0.46))
        icon.name = name
        icon.fill.background()
        icon.line.color.rgb = green
        icon.line.width = Pt(2)
        check = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(cx - 0.09), inch(y + 0.22), inch(cx - 0.01), inch(y + 0.31))
        check.line.color.rgb = green
        check.line.width = Pt(2)
        check2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(cx - 0.01), inch(y + 0.31), inch(cx + 0.13), inch(y + 0.14))
        check2.line.color.rgb = green
        check2.line.width = Pt(2)
    elif index == 3:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CLOUD, inch(cx - 0.26), inch(y + 0.04), inch(0.52), inch(0.34))
        icon.name = name
        icon.fill.background()
        icon.line.color.rgb = green
        icon.line.width = Pt(2)
        arm = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(cx + 0.06), inch(y + 0.26), inch(cx + 0.22), inch(y + 0.14))
        arm.line.color.rgb = green
        arm.line.width = Pt(2.5)
    elif index == 4:
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - 0.05), inch(y), inch(0.10), inch(0.10))
        head.name = name
        set_fill(head, WHITE)
        head.line.color.rgb = green
        head.line.width = Pt(1.8)
        for x1, y1, x2, y2 in [
            (-0.02, 0.10, -0.12, 0.26),
            (-0.02, 0.10, 0.13, 0.22),
            (-0.10, 0.26, -0.22, 0.43),
            (-0.10, 0.26, 0.04, 0.43),
            (0.08, 0.19, 0.24, 0.12),
        ]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(cx + x1), inch(y + y1), inch(cx + x2), inch(y + y2))
            line.line.color.rgb = green
            line.line.width = Pt(2)
    elif index == 5:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CLOUD, inch(cx - 0.24), inch(y + 0.02), inch(0.48), inch(0.42))
        icon.name = name
        icon.fill.background()
        icon.line.color.rgb = green
        icon.line.width = Pt(2)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(cx), inch(y + 0.04), inch(cx), inch(y + 0.42))
        line.line.color.rgb = green
        line.line.width = Pt(1.2)
    else:
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - 0.12), inch(y), inch(0.24), inch(0.24))
        head.name = name
        head.fill.background()
        head.line.color.rgb = green
        head.line.width = Pt(2)
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ARC, inch(cx - 0.25), inch(y + 0.28), inch(0.5), inch(0.42))
        body.fill.background()
        body.line.color.rgb = green
        body.line.width = Pt(2)


def add_vertical_dashed_line(slide, x, y1, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x), inch(y1), inch(x), inch(y2))
    line.line.color.rgb = LIGHT_GREEN
    line.line.width = Pt(1)
    line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_contact_icon(slide, kind, x, y):
    if kind == "pin":
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x), inch(y), inch(0.18), inch(0.18))
        icon.name = "ICON_ADDRESS"
        set_fill(icon, WHITE)
        icon.line.color.rgb = GREEN_DARK
        icon.line.width = Pt(2)
        tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, inch(x + 0.045), inch(y + 0.12), inch(0.09), inch(0.13))
        set_fill(tri, GREEN_DARK)
        tri.line.color.rgb = GREEN_DARK
    elif kind == "phone":
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(0.17), inch(0.24))
        icon.name = "ICON_PHONE"
        set_fill(icon, GREEN_DARK)
        icon.line.color.rgb = GREEN_DARK
    else:
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x), inch(y), inch(0.2), inch(0.2))
        icon.name = "ICON_INSTAGRAM"
        icon.fill.background()
        icon.line.color.rgb = GREEN_DARK
        icon.line.width = Pt(1.8)


def add_brush_cta(slide):
    base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(0), inch(10.28), inch(3.7), inch(1.41))
    base.name = "CTA_BRUSH_BASE"
    set_fill(base, BLACK)
    set_no_line(base)
    for i, x in enumerate([3.42, 3.52, 3.60, 3.66]):
        rag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, inch(x), inch(10.25 + i * 0.04), inch(0.42), inch(1.26 - i * 0.12))
        rag.name = f"CTA_BRUSH_RAGGED_{i + 1}"
        set_fill(rag, BLACK)
        set_no_line(rag)
    for i in range(9):
        strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(0.0), inch(10.22 + i * 0.16), inch(3.48 - i * 0.18), inch(0.035))
        strip.name = f"CTA_BRUSH_STROKE_{i + 1}"
        set_fill(strip, BLACK)
        set_no_line(strip)

    main = add_text(slide, 0.28, 10.42, 2.95, 0.32, "{{CTA_MAIN}}", 22, GREEN, DISPLAY_FALLBACK, True, PP_ALIGN.LEFT, "TEXT_CTA_MAIN")
    sub = add_text(slide, 0.28, 10.75, 3.05, 0.72, "{{CTA_SUB}}", 22, WHITE, DISPLAY_FALLBACK, True, PP_ALIGN.LEFT, "TEXT_CTA_SUB")
    main.text_frame.word_wrap = False
    sub.text_frame.word_wrap = True

    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, inch(3.06), inch(10.39), inch(0.75), inch(0.94))
    arrow.name = "CTA_GREEN_ARROW"
    set_fill(arrow, GREEN)
    set_no_line(arrow)


def build_template():
    TEMPLATE_PATH.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Hero image and high-energy green strokes.
    add_placeholder(slide, "IMAGE_HERO", 3.42, 0.0, 4.85, 5.08, "IMAGE_HERO", RGBColor(236, 238, 235))
    add_skew_bar(slide, 3.19, 0.0, 0.06, 1.45, GREEN, -22, "HERO_GREEN_STROKE_1")
    add_skew_bar(slide, 3.06, 0.18, 0.04, 0.95, GREEN, -22, "HERO_GREEN_STROKE_2")
    add_skew_bar(slide, 3.25, 0.34, 0.03, 0.75, GREEN, -22, "HERO_GREEN_STROKE_3")

    # Top logo block.
    add_logo_mark_placeholder(slide, "LOGO_TOP", 0.38, 0.34, 0.72)
    add_text(slide, 1.28, 0.45, 1.7, 0.26, "VIRTUS", 24, BLACK, BODY_FONT, True, PP_ALIGN.LEFT, "BRAND_TOP_NAME")
    add_text(slide, 1.30, 0.75, 1.85, 0.18, "COMBAT ACADEMY", 8.5, BLACK, BODY_FONT, True, PP_ALIGN.LEFT, "BRAND_TOP_SUB")

    # Main headline.
    add_text(slide, 0.13, 1.28, 3.3, 0.86, "{{TITLE_1}}", 60, GREEN_DARK, DISPLAY_FALLBACK, True, PP_ALIGN.LEFT, "TEXT_TITLE_1", -4)
    add_text(slide, 0.18, 2.1, 3.85, 1.06, "{{TITLE_2}}", 58, BLACK, DISPLAY_FALLBACK, True, PP_ALIGN.LEFT, "TEXT_TITLE_2", -4)
    add_age_box(slide)

    # Gallery strip.
    gallery_y = 5.08
    gallery_h = 2.0
    add_placeholder(slide, "IMAGE_GALLERY_1", 0.0, gallery_y, 2.72, gallery_h, "IMAGE_GALLERY_1", RGBColor(232, 235, 230))
    add_placeholder(slide, "IMAGE_GALLERY_2", 2.72, gallery_y, 2.82, gallery_h, "IMAGE_GALLERY_2", RGBColor(232, 235, 230))
    add_placeholder(slide, "IMAGE_GALLERY_3", 5.54, gallery_y, 2.73, gallery_h, "IMAGE_GALLERY_3", RGBColor(232, 235, 230))
    add_skew_bar(slide, 2.67, gallery_y - 0.03, 0.045, gallery_h + 0.09, GREEN, 8, "GALLERY_SEPARATOR_1")
    add_skew_bar(slide, 5.49, gallery_y - 0.03, 0.045, gallery_h + 0.09, GREEN, 8, "GALLERY_SEPARATOR_2")

    # Benefits title and columns.
    title_y = 7.22
    left_rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(0.35), inch(title_y + 0.20), inch(2.68), inch(title_y + 0.20))
    left_rule.name = "BENEFITS_RULE_LEFT"
    left_rule.line.color.rgb = GREEN_DARK
    left_rule.line.width = Pt(1)
    right_rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(5.65), inch(title_y + 0.20), inch(8.02), inch(title_y + 0.20))
    right_rule.name = "BENEFITS_RULE_RIGHT"
    right_rule.line.color.rgb = GREEN_DARK
    right_rule.line.width = Pt(1)
    add_text(slide, 2.78, title_y, 2.75, 0.42, "{{BENEFITS_TITLE}}", 24, BLACK, DISPLAY_FALLBACK, True, PP_ALIGN.CENTER, "TEXT_BENEFITS_TITLE")

    columns = [
        ("{{BENEFIT_1_TITLE}}", "{{BENEFIT_1_TEXT}}"),
        ("{{BENEFIT_2_TITLE}}", "{{BENEFIT_2_TEXT}}"),
        ("{{BENEFIT_3_TITLE}}", "{{BENEFIT_3_TEXT}}"),
        ("{{BENEFIT_4_TITLE}}", "{{BENEFIT_4_TEXT}}"),
        ("{{BENEFIT_5_TITLE}}", "{{BENEFIT_5_TEXT}}"),
        ("{{BENEFIT_6_TITLE}}", "{{BENEFIT_6_TEXT}}"),
    ]
    col_w = 1.31
    x0 = 0.22
    for i, (title, body) in enumerate(columns, start=1):
        x = x0 + (i - 1) * col_w
        cx = x + col_w / 2
        add_benefit_icon(slide, i, cx, 7.75)
        add_text(slide, x + 0.05, 8.39, col_w - 0.1, 0.22, title, 13.5, BLACK, DISPLAY_FALLBACK, True, PP_ALIGN.CENTER, f"TEXT_BENEFIT_{i}_TITLE")
        body_box = add_text(slide, x + 0.11, 8.78, col_w - 0.22, 0.92, body, 8.1, TEXT_GRAY, BODY_FONT, False, PP_ALIGN.CENTER, f"TEXT_BENEFIT_{i}_TEXT")
        body_box.text_frame.word_wrap = True
        if i < 6:
            add_vertical_dashed_line(slide, x + col_w, 7.72, 9.78).name = f"BENEFIT_SEPARATOR_{i}"

    # Bottom CTA and contact footer.
    add_brush_cta(slide)
    contact_x = 4.08
    add_contact_icon(slide, "pin", contact_x, 10.37)
    add_text(slide, contact_x + 0.33, 10.36, 1.95, 0.22, "{{ADDRESS}}", 11.2, BLACK, BODY_FONT, False, PP_ALIGN.LEFT, "TEXT_ADDRESS")
    add_contact_icon(slide, "phone", contact_x, 10.78)
    add_text(slide, contact_x + 0.33, 10.77, 1.95, 0.22, "{{PHONE}}", 11.2, BLACK, BODY_FONT, False, PP_ALIGN.LEFT, "TEXT_PHONE")
    add_contact_icon(slide, "instagram", contact_x, 11.18)
    add_text(slide, contact_x + 0.33, 11.17, 2.05, 0.22, "{{INSTAGRAM}}", 11.2, BLACK, BODY_FONT, False, PP_ALIGN.LEFT, "TEXT_INSTAGRAM")

    add_logo_mark_placeholder(slide, "LOGO_BOTTOM", 6.42, 10.27, 0.62)
    add_text(slide, 7.04, 10.37, 0.95, 0.20, "VIRTUS", 17, BLACK, BODY_FONT, True, PP_ALIGN.LEFT, "BRAND_BOTTOM_NAME")
    add_text(slide, 7.05, 10.60, 1.0, 0.16, "COMBAT ACADEMY", 5.8, BLACK, BODY_FONT, True, PP_ALIGN.LEFT, "BRAND_BOTTOM_SUB")

    instagram = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inch(6.94), inch(11.08), inch(0.26), inch(0.26))
    instagram.name = "SOCIAL_INSTAGRAM_ICON"
    instagram.fill.background()
    instagram.line.color.rgb = BLACK
    instagram.line.width = Pt(1.4)
    facebook = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(7.62), inch(11.08), inch(0.28), inch(0.28))
    facebook.name = "SOCIAL_FACEBOOK_ICON"
    facebook.fill.background()
    facebook.line.color.rgb = BLACK
    facebook.line.width = Pt(1.4)
    add_text(slide, 7.69, 11.105, 0.13, 0.18, "f", 11, BLACK, BODY_FONT, True, PP_ALIGN.CENTER, "SOCIAL_FACEBOOK_F")

    prs.save(TEMPLATE_PATH)
    return TEMPLATE_PATH


if __name__ == "__main__":
    output = build_template()
    print(f"Template creato: {output}")
