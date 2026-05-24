from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt


DEFAULT_SLIDE_WIDTH_IN = 7.5


def render_layout_to_pptx(layout, project_root, output_path):
    project_root = Path(project_root)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    canvas = layout["canvas"]
    canvas_width = float(canvas["width"])
    canvas_height = float(canvas["height"])
    scale = DEFAULT_SLIDE_WIDTH_IN / canvas_width
    slide_width = DEFAULT_SLIDE_WIDTH_IN
    slide_height = canvas_height * scale

    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = canvas.get("background", "#FFFFFF")
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = parse_color(background)

    assets_base = resolve_assets_base(layout.get("assets_base", ""), project_root)
    warnings = []

    for layer in layout.get("layers", []):
        layer_type = layer["type"]
        if layer_type == "rect":
            add_rect(slide, layer, scale)
        elif layer_type == "text":
            add_text(slide, layer, scale)
        elif layer_type == "image":
            add_image(slide, layer, scale, assets_base, warnings)
        elif layer_type == "line":
            add_line(slide, layer, scale)

    prs.save(output_path)
    return output_path, warnings


def resolve_assets_base(value, project_root):
    if not value:
        return project_root
    path = Path(value)
    if path.is_absolute():
        return path
    return project_root / path


def px(value, scale):
    return Inches(float(value) * scale)


def parse_color(value):
    if isinstance(value, RGBColor):
        return value
    text = str(value).strip().lstrip("#")
    if len(text) != 6:
        raise ValueError(f"Colore non valido: {value}")
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def opacity_to_transparency(opacity):
    if opacity is None:
        return 0
    opacity = max(0.0, min(1.0, float(opacity)))
    return int(round((1.0 - opacity) * 100))


def add_rect(slide, layer, scale):
    radius = float(layer.get("radius", 0) or 0)
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius > 0 else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        px(layer["x"], scale),
        px(layer["y"], scale),
        px(layer["w"], scale),
        px(layer["h"], scale),
    )
    shape.name = layer.get("id", "rect")
    shape.fill.solid()
    shape.fill.fore_color.rgb = parse_color(layer.get("fill", "#FFFFFF"))
    shape.fill.transparency = opacity_to_transparency(layer.get("opacity"))
    shape.line.fill.background()
    return shape


def add_text(slide, layer, scale):
    box = slide.shapes.add_textbox(
        px(layer["x"], scale),
        px(layer["y"], scale),
        px(layer["w"], scale),
        px(layer["h"], scale),
    )
    box.name = layer.get("id", "text")
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP

    text = str(layer.get("text", ""))
    if layer.get("uppercase"):
        text = text.upper()

    paragraph = frame.paragraphs[0]
    paragraph.alignment = parse_align(layer.get("align", "left"))
    run = paragraph.add_run()
    run.text = text
    run.font.name = layer.get("font", "Arial")
    run.font.size = Pt(float(layer.get("size", 24)))
    run.font.bold = bool(layer.get("bold", False))
    run.font.italic = bool(layer.get("italic", False))
    run.font.color.rgb = parse_color(layer.get("color", "#000000"))
    return box


def parse_align(value):
    value = str(value).lower()
    if value == "center":
        return PP_ALIGN.CENTER
    if value == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def add_image(slide, layer, scale, assets_base, warnings):
    image_path = assets_base / layer["src"]
    if not image_path.exists():
        message = f"Immagine mancante per layer '{layer.get('id', 'image')}': {image_path}"
        warnings.append(message)
        print(f"[WARNING] {message}")
        return None

    x = px(layer["x"], scale)
    y = px(layer["y"], scale)
    w = px(layer["w"], scale)
    h = px(layer["h"], scale)
    picture = slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)
    picture.name = layer.get("id", "image")
    picture.crop_left = 0
    picture.crop_right = 0
    picture.crop_top = 0
    picture.crop_bottom = 0

    fit = str(layer.get("fit", "cover")).lower()
    if fit == "cover":
        apply_cover_crop(picture, image_path, float(layer["w"]), float(layer["h"]))
    elif fit == "contain":
        apply_contain_size(picture, image_path, x, y, w, h)

    if "opacity" in layer:
        apply_picture_opacity(picture, layer["opacity"])
    return picture


def apply_picture_opacity(picture, opacity):
    opacity = max(0.0, min(1.0, float(opacity)))
    alpha = int(round(opacity * 100000))
    blip = picture._element.blipFill.blip
    for child in list(blip):
        if child.tag.endswith("alphaModFix"):
            blip.remove(child)
    blip.append(parse_xml(f'<a:alphaModFix {nsdecls("a")} amt="{alpha}"/>'))


def apply_cover_crop(picture, image_path, box_w, box_h):
    with Image.open(image_path) as image:
        image_w, image_h = image.size
    image_aspect = image_w / image_h
    box_aspect = box_w / box_h

    if image_aspect > box_aspect:
        crop = (1 - box_aspect / image_aspect) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_aspect < box_aspect:
        crop = (1 - image_aspect / box_aspect) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop


def apply_contain_size(picture, image_path, x, y, box_w, box_h):
    with Image.open(image_path) as image:
        image_w, image_h = image.size
    image_aspect = image_w / image_h
    box_aspect = box_w / box_h

    if image_aspect > box_aspect:
        new_w = box_w
        new_h = int(box_w / image_aspect)
    else:
        new_h = box_h
        new_w = int(box_h * image_aspect)

    picture.width = new_w
    picture.height = new_h
    picture.left = x + int((box_w - new_w) / 2)
    picture.top = y + int((box_h - new_h) / 2)


def add_line(slide, layer, scale):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px(layer["x1"], scale),
        px(layer["y1"], scale),
        px(layer["x2"], scale),
        px(layer["y2"], scale),
    )
    line.name = layer.get("id", "line")
    line.line.color.rgb = parse_color(layer.get("color", "#000000"))
    line.line.width = Pt(float(layer.get("width", 1)))
    return line
