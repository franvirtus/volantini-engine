import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


TEXT_PLACEHOLDER_RE = re.compile(r"\{\{([A-Za-z0-9_]+)\}\}")
IMAGE_TOKEN_RE = re.compile(r"\bIMAGE_[A-Za-z0-9_]+\b")


class RenderReport:
    def __init__(self):
        self.text_found = []
        self.text_replaced = []
        self.missing_text_keys = set()
        self.images_found = []
        self.images_replaced = []
        self.missing_image_placeholders = set()
        self.missing_image_files = []
        self.remaining_placeholders = []
        self.warnings = []


def warn(report, message):
    report.warnings.append(message)
    print(f"[WARNING] {message}")


def iter_shapes(shapes):
    for shape in list(shapes):
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)


def collect_text_placeholders(prs):
    found = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes):
            for key in TEXT_PLACEHOLDER_RE.findall(shape_text(shape)):
                found.append(
                    {
                        "slide": slide_number,
                        "shape": getattr(shape, "name", ""),
                        "key": key,
                        "placeholder": "{{" + key + "}}",
                    }
                )
    return found


def collect_image_placeholders(prs, image_keys):
    image_keys = set(image_keys)
    found = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes):
            shape_name = getattr(shape, "name", "")
            if shape_name in image_keys:
                found.append({"slide": slide_number, "shape": shape_name})
    return found


def replace_placeholders_in_text(text, texts, report):
    replacements = []

    def replace_match(match):
        key = match.group(1)
        placeholder = match.group(0)
        if key not in texts:
            report.missing_text_keys.add(key)
            return placeholder
        replacements.append(placeholder)
        value = texts[key]
        return "" if value is None else str(value)

    return TEXT_PLACEHOLDER_RE.sub(replace_match, text), replacements


def replace_text_in_paragraph(paragraph, texts, report):
    original = paragraph.text
    if not original:
        return []

    updated, replacements = replace_placeholders_in_text(original, texts, report)
    if updated == original:
        return []

    # Rewriting the first run preserves the placeholder box style and also works
    # when PowerPoint has split one placeholder across several runs.
    if paragraph.runs:
        paragraph.runs[0].text = updated
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = updated

    return replacements


def replace_all_text(prs, texts, report):
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_count = 0
        for shape in iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            for paragraph in shape.text_frame.paragraphs:
                replacements = replace_text_in_paragraph(paragraph, texts, report)
                for placeholder in replacements:
                    report.text_replaced.append(
                        {
                            "slide": slide_number,
                            "shape": getattr(shape, "name", ""),
                            "placeholder": placeholder,
                        }
                    )
                slide_count += len(replacements)

        print(f"[VolantiniEngine] Slide {slide_number}: sostituzioni testo eseguite: {slide_count}")


def replacement_picture_name(placeholder_name):
    if placeholder_name.startswith("IMAGE_"):
        return "PICTURE_" + placeholder_name.removeprefix("IMAGE_")
    if placeholder_name.startswith("LOGO_"):
        return "PICTURE_" + placeholder_name
    return "PICTURE_" + placeholder_name


def add_picture_in_same_position(slide, image_path, reference_shape):
    new_picture = slide.shapes.add_picture(
        str(image_path),
        reference_shape.left,
        reference_shape.top,
        width=reference_shape.width,
        height=reference_shape.height,
    )
    new_picture.name = replacement_picture_name(reference_shape.name)

    reference_element = reference_shape._element
    new_picture_element = new_picture._element
    reference_parent = reference_element.getparent()
    reference_index = reference_parent.index(reference_element)

    reference_parent.remove(new_picture_element)
    reference_parent.insert(reference_index, new_picture_element)
    reference_parent.remove(reference_element)
    return new_picture


def replace_all_images(prs, images, project_root, report):
    image_keys = set(images.keys())
    found_keys = set()

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_count = 0
        for shape in list(iter_shapes(slide.shapes)):
            shape_name = getattr(shape, "name", "")
            if shape_name not in image_keys:
                continue

            found_keys.add(shape_name)
            image_path = resolve_project_path(images[shape_name], project_root)
            if not image_path.exists():
                report.missing_image_files.append({"shape": shape_name, "path": str(image_path)})
                warn(report, f"File immagine mancante per '{shape_name}': {image_path}. Placeholder lasciato invariato.")
                continue

            add_picture_in_same_position(slide, image_path, shape)
            report.images_replaced.append(
                {"slide": slide_number, "shape": shape_name, "path": str(image_path)}
            )
            slide_count += 1

        print(f"[VolantiniEngine] Slide {slide_number}: immagini sostituite: {slide_count}")

    for key in sorted(image_keys - found_keys):
        report.missing_image_placeholders.add(key)
        warn(report, f"Placeholder immagine '{key}' non trovato nel template.")


def collect_remaining_placeholders(prs):
    remaining = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes):
            name = getattr(shape, "name", "")
            text = shape_text(shape)
            for key in TEXT_PLACEHOLDER_RE.findall(text):
                remaining.append(
                    {"type": "text", "slide": slide_number, "shape": name, "value": "{{" + key + "}}"}
                )
            for image_token in IMAGE_TOKEN_RE.findall(text):
                remaining.append(
                    {"type": "image-text", "slide": slide_number, "shape": name, "value": image_token}
                )
            if name.startswith("IMAGE_"):
                remaining.append(
                    {"type": "image-shape", "slide": slide_number, "shape": name, "value": name}
                )
    return remaining


def resolve_project_path(value, project_root):
    path = Path(value)
    if path.is_absolute():
        return path
    return project_root / path


def render_pptx(config, project_root, out_dir):
    project_root = Path(project_root)
    out_dir = Path(out_dir)
    template_path = project_root / "templates" / config["template"]

    if not template_path.exists():
        raise FileNotFoundError(f"Template PowerPoint non trovato: {template_path}")

    out_dir.mkdir(parents=True, exist_ok=True)
    output_path = out_dir / f"{config['output_name']}.pptx"

    print(f"[VolantiniEngine] Apro template: {template_path}")
    prs = Presentation(str(template_path))
    report = RenderReport()

    texts = config.get("texts", {})
    images = config.get("images", {})
    report.text_found = collect_text_placeholders(prs)
    report.images_found = collect_image_placeholders(prs, images.keys())

    print(f"[VolantiniEngine] Placeholder testo trovati nel template: {len(report.text_found)}")
    print(f"[VolantiniEngine] Placeholder testo nel config: {len(texts)}")
    replace_all_text(prs, texts, report)

    for key in sorted(report.missing_text_keys):
        warn(report, f"Chiave testo mancante nel config: '{key}'")

    print(f"[VolantiniEngine] Placeholder immagini trovati nel template: {len(report.images_found)}")
    print(f"[VolantiniEngine] Placeholder immagini nel config: {len(images)}")
    replace_all_images(prs, images, project_root, report)

    report.remaining_placeholders = collect_remaining_placeholders(prs)
    for item in report.remaining_placeholders:
        warn(
            report,
            f"Placeholder rimasto ({item['type']}): slide {item['slide']} | {item['shape']} | {item['value']}",
        )

    prs.save(str(output_path))
    print(f"[VolantiniEngine] PowerPoint generato: {output_path}")
    print(f"[VolantiniEngine] Totale sostituzioni testo: {len(report.text_replaced)}")
    print(f"[VolantiniEngine] Totale immagini sostituite: {len(report.images_replaced)}")
    return output_path, report

