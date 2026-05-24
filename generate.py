import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PROJECT_ROOT = Path(__file__).resolve().parent
TEMPLATES_DIR = PROJECT_ROOT / "templates"
OUTPUT_DIR = PROJECT_ROOT / "output"
POWERPOINT_FORMAT_PDF = 32
POWERPOINT_FORMAT_PNG = 18
TEXT_PLACEHOLDER_RE = re.compile(r"\{\{([A-Za-z0-9_]+)\}\}")
IMAGE_TOKEN_RE = re.compile(r"\bIMAGE_[A-Za-z0-9_]+\b")


def log(message):
    print(f"[VolantiniEngine] {message}")


def warn(message):
    print(f"[WARNING] {message}")


def load_config(config_path):
    try:
        with config_path.open("r", encoding="utf-8") as file:
            return json.load(file)
    except FileNotFoundError:
        raise FileNotFoundError(f"File JSON non trovato: {config_path}")
    except json.JSONDecodeError as exc:
        raise ValueError(f"JSON non valido in {config_path}: {exc}") from exc


def resolve_project_path(value):
    path = Path(value)
    if path.is_absolute():
        return path
    return PROJECT_ROOT / path


def iter_shapes(shapes):
    for shape in list(shapes):
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)


def collect_text_placeholders(prs):
    found = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes):
            for key in TEXT_PLACEHOLDER_RE.findall(shape_text(shape)):
                found.append(
                    {
                        "slide": slide_number,
                        "shape": getattr(shape, "name", ""),
                        "key": key,
                        "placeholder": "{{" + key + "}}",
                    }
                )
    return found


def collect_image_placeholders(prs, image_keys):
    found = []
    image_keys = set(image_keys)
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes):
            shape_name = getattr(shape, "name", "")
            if shape_name in image_keys:
                found.append({"slide": slide_number, "shape": shape_name})
    return found


def replace_placeholders_in_text(text, texts, missing_keys):
    replacements = []

    def replace_match(match):
        key = match.group(1)
        placeholder = match.group(0)
        if key not in texts:
            missing_keys.add(key)
            return placeholder
        replacements.append(placeholder)
        value = texts[key]
        return "" if value is None else str(value)

    return TEXT_PLACEHOLDER_RE.sub(replace_match, text), replacements


def replace_text_in_paragraph(paragraph, texts, missing_keys):
    original = paragraph.text
    if not original:
        return []

    updated, replacements = replace_placeholders_in_text(original, texts, missing_keys)
    if updated == original:
        return []

    if paragraph.runs:
        paragraph.runs[0].text = updated
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = updated

    return replacements


def replace_text_in_shape(shape, texts, missing_keys):
    if not getattr(shape, "has_text_frame", False):
        return []

    replacements = []
    for paragraph in shape.text_frame.paragraphs:
        replacements.extend(replace_text_in_paragraph(paragraph, texts, missing_keys))
    return replacements


def replace_all_text(prs, texts):
    missing_keys = set()
    replacements = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_count = 0
        for shape in iter_shapes(slide.shapes):
            shape_replacements = replace_text_in_shape(shape, texts, missing_keys)
            for placeholder in shape_replacements:
                replacements.append(
                    {
                        "slide": slide_number,
                        "shape": getattr(shape, "name", ""),
                        "placeholder": placeholder,
                    }
                )
            slide_count += len(shape_replacements)
        log(f"Slide {slide_number}: sostituzioni testo eseguite: {slide_count}")

    for key in sorted(missing_keys):
        warn(f"Chiave testo mancante nel JSON: '{key}'")

    return replacements, missing_keys


def replacement_picture_name(placeholder_name):
    if placeholder_name.startswith("IMAGE_"):
        return "PICTURE_" + placeholder_name.removeprefix("IMAGE_")
    if placeholder_name.startswith("LOGO_"):
        return "PICTURE_" + placeholder_name
    return "PICTURE_" + placeholder_name


def add_picture_in_same_position(slide, image_path, reference_shape):
    new_picture = slide.shapes.add_picture(
        str(image_path),
        reference_shape.left,
        reference_shape.top,
        width=reference_shape.width,
        height=reference_shape.height,
    )

    new_picture.name = replacement_picture_name(reference_shape.name)

    reference_element = reference_shape._element
    new_picture_element = new_picture._element
    reference_parent = reference_element.getparent()
    reference_index = reference_parent.index(reference_element)

    reference_parent.remove(new_picture_element)
    reference_parent.insert(reference_index, new_picture_element)
    reference_parent.remove(reference_element)

    return new_picture


def replace_image_shape(slide, shape, image_path):
    if not image_path.exists():
        warn(f"File immagine mancante per '{shape.name}': {image_path}. Placeholder lasciato invariato.")
        return False

    add_picture_in_same_position(slide, image_path, shape)
    return True


def replace_all_images(prs, images):
    replacements = []
    image_keys = set(images.keys())
    found_keys = set()

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_count = 0
        for shape in list(iter_shapes(slide.shapes)):
            shape_name = getattr(shape, "name", "")
            if shape_name not in image_keys:
                continue

            found_keys.add(shape_name)
            image_path = resolve_project_path(images[shape_name])
            if replace_image_shape(slide, shape, image_path):
                replacements.append(
                    {
                        "slide": slide_number,
                        "shape": shape_name,
                        "path": str(image_path),
                    }
                )
                slide_count += 1

        log(f"Slide {slide_number}: immagini sostituite: {slide_count}")

    for key in sorted(image_keys - found_keys):
        warn(f"Placeholder immagine '{key}' non trovato nel template.")

    return replacements, image_keys - found_keys


def validate_config(config):
    required_fields = ["template", "output_name"]
    missing = [field for field in required_fields if not config.get(field)]

    if missing:
        raise ValueError("Campi obbligatori mancanti nel JSON: " + ", ".join(missing))

    if "texts" in config and not isinstance(config["texts"], dict):
        raise ValueError("Il campo 'texts' deve essere un oggetto JSON.")

    if "images" in config and not isinstance(config["images"], dict):
        raise ValueError("Il campo 'images' deve essere un oggetto JSON.")


def remaining_placeholders(prs):
    remaining = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes):
            name = getattr(shape, "name", "")
            text = shape_text(shape)
            for placeholder in TEXT_PLACEHOLDER_RE.findall(text):
                remaining.append(
                    {
                        "type": "text",
                        "slide": slide_number,
                        "shape": name,
                        "value": "{{" + placeholder + "}}",
                    }
                )
            for image_token in IMAGE_TOKEN_RE.findall(text):
                remaining.append(
                    {
                        "type": "image-text",
                        "slide": slide_number,
                        "shape": name,
                        "value": image_token,
                    }
                )
            if name.startswith("IMAGE_"):
                remaining.append(
                    {
                        "type": "image-shape",
                        "slide": slide_number,
                        "shape": name,
                        "value": name,
                    }
                )
    return remaining


def print_debug_report(report):
    print("\n[DEBUG] Template usato:")
    print(f"  {report['template_path']}")
    print("[DEBUG] Output creato:")
    print(f"  {report['output_path']}")

    print("[DEBUG] Placeholder testo trovati:")
    if report["text_found"]:
        for item in report["text_found"]:
            print(f"  slide {item['slide']} | {item['shape']} | {item['placeholder']}")
    else:
        print("  Nessuno")

    print("[DEBUG] Placeholder testo sostituiti:")
    if report["text_replaced"]:
        for item in report["text_replaced"]:
            print(f"  slide {item['slide']} | {item['shape']} | {item['placeholder']}")
    else:
        print("  Nessuno")

    print("[DEBUG] Immagini trovate:")
    if report["images_found"]:
        for item in report["images_found"]:
            print(f"  slide {item['slide']} | {item['shape']}")
    else:
        print("  Nessuna")

    print("[DEBUG] Immagini sostituite:")
    if report["images_replaced"]:
        for item in report["images_replaced"]:
            print(f"  slide {item['slide']} | {item['shape']} -> {item['path']}")
    else:
        print("  Nessuna")

    print("[DEBUG] Placeholder rimasti:")
    if report["remaining"]:
        for item in report["remaining"]:
            print(f"  {item['type']} | slide {item['slide']} | {item['shape']} | {item['value']}")
    else:
        print("  Nessuno")


def build_presentation(config_path, debug=False):
    log(f"Leggo configurazione: {config_path}")
    config = load_config(config_path)
    validate_config(config)

    template_path = TEMPLATES_DIR / config["template"]
    if not template_path.exists():
        raise FileNotFoundError(f"Template PowerPoint non trovato: {template_path}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUTPUT_DIR / f"{config['output_name']}.pptx"

    log(f"Apro template: {template_path}")
    prs = Presentation(str(template_path))

    texts = config.get("texts", {})
    images = config.get("images", {})
    text_found = collect_text_placeholders(prs)
    images_found = collect_image_placeholders(prs, images.keys())

    log(f"Placeholder testo trovati nel template: {len(text_found)}")
    log(f"Placeholder testo nel JSON: {len(texts)}")
    text_replacements, missing_text_keys = replace_all_text(prs, texts)

    log(f"Placeholder immagini trovati nel template: {len(images_found)}")
    log(f"Placeholder immagini nel JSON: {len(images)}")
    image_replacements, missing_image_placeholders = replace_all_images(prs, images)

    remaining = remaining_placeholders(prs)
    for item in remaining:
        warn(f"Placeholder rimasto ({item['type']}): slide {item['slide']} | {item['shape']} | {item['value']}")

    prs.save(str(output_path))

    log(f"PowerPoint generato: {output_path}")
    log(f"Totale sostituzioni testo: {len(text_replacements)}")
    log(f"Totale immagini sostituite: {len(image_replacements)}")

    report = {
        "template_path": str(template_path),
        "output_path": str(output_path),
        "text_found": text_found,
        "text_replaced": text_replacements,
        "missing_text_keys": sorted(missing_text_keys),
        "images_found": images_found,
        "images_replaced": image_replacements,
        "missing_image_placeholders": sorted(missing_image_placeholders),
        "remaining": remaining,
    }

    if debug:
        print_debug_report(report)

    return output_path, report


def export_with_powerpoint(pptx_path, export_pdf=False, export_png=False):
    if not export_pdf and not export_png:
        return

    try:
        import win32com.client
    except ImportError:
        warn(
            "Esportazione PDF/PNG richiesta ma pywin32 non e installato. "
            "Il PowerPoint modificabile e stato comunque generato."
        )
        warn("Per abilitare export automatico: python -m pip install pywin32")
        return

    powerpoint = None
    presentation = None

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)

        if export_pdf:
            pdf_path = pptx_path.with_suffix(".pdf")
            presentation.SaveAs(str(pdf_path), POWERPOINT_FORMAT_PDF)
            log(f"PDF generato: {pdf_path}")

        if export_png:
            png_dir = pptx_path.with_suffix("")
            presentation.SaveAs(str(png_dir), POWERPOINT_FORMAT_PNG)
            log(f"PNG generati nella cartella: {png_dir}")

    except Exception as exc:
        warn(f"Export PDF/PNG non completato: {exc}")
    finally:
        if presentation is not None:
            presentation.Close()
        if powerpoint is not None:
            powerpoint.Quit()


def parse_args():
    parser = argparse.ArgumentParser(
        description="Generatore generico di volantini PowerPoint modificabili per Virtus."
    )
    parser.add_argument(
        "json_file",
        help="Percorso del file JSON contenuti, per esempio data/examples/boxe_10_14.json",
    )
    parser.add_argument(
        "--debug",
        action="store_true",
        help="Stampa template, output, placeholder trovati/sostituiti e placeholder rimasti.",
    )
    parser.add_argument(
        "--pdf",
        action="store_true",
        help="Esporta anche PDF usando PowerPoint installato su Windows e pywin32.",
    )
    parser.add_argument(
        "--png",
        action="store_true",
        help="Esporta anche PNG usando PowerPoint installato su Windows e pywin32.",
    )
    return parser.parse_args()


def main():
    args = parse_args()
    config_path = resolve_project_path(args.json_file)

    try:
        output_path, _report = build_presentation(config_path, debug=args.debug)
        export_with_powerpoint(output_path, export_pdf=args.pdf, export_png=args.png)
    except Exception as exc:
        print(f"[ERRORE] {exc}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
