from pathlib import Path
from math import sin, cos, pi

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "assets" / "images"
TEMPLATE_DIR = ROOT / "templates"
DATA_DIR = ROOT / "data" / "examples"

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(12, 13, 12)
ANTHRACITE = RGBColor(38, 42, 38)
MUTED = RGBColor(98, 105, 98)
LIGHT = RGBColor(242, 244, 240)
GREEN = RGBColor(183, 255, 0)
GREEN_DARK = RGBColor(116, 170, 0)


def ensure_dirs():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)
    DATA_DIR.mkdir(parents=True, exist_ok=True)


def make_photo(path, variant):
    width, height = 1600, 1100
    img = Image.new("RGB", (width, height), "#eef2eb")
    draw = ImageDraw.Draw(img, "RGBA")

    for y in range(height):
        tone = int(244 - y * 26 / height)
        draw.line([(0, y), (width, y)], fill=(tone, tone + 3, tone - 2, 255))

    # gym wall, mats, windows
    draw.rectangle([0, int(height * 0.58), width, height], fill=(218, 224, 214, 255))
    draw.rectangle([0, int(height * 0.6), width, int(height * 0.64)], fill=(185, 255, 0, 150))
    for x in range(70, width, 310):
        draw.rounded_rectangle([x, 80, x + 210, 330], radius=20, fill=(255, 255, 255, 180))
        draw.rectangle([x + 12, 92, x + 198, 318], outline=(205, 214, 203, 160), width=4)

    if variant == "hero":
        people = [
            (420, 610, "#1e211f", "#b7ff00"),
            (620, 610, "#2a2e2a", "#111111"),
            (820, 610, "#1c1d1c", "#b7ff00"),
            (1030, 620, "#303530", "#111111"),
        ]
    elif variant == "drill":
        people = [(580, 635, "#202320", "#b7ff00"), (880, 635, "#262b26", "#111111")]
    elif variant == "wrestling":
        people = [(610, 665, "#202420", "#111111"), (900, 665, "#252925", "#b7ff00")]
    else:
        people = [(750, 635, "#202320", "#b7ff00")]
        draw.ellipse([1040, 260, 1190, 730], fill=(28, 30, 28, 255))
        draw.rectangle([1100, 170, 1130, 270], fill=(35, 38, 35, 255))

    for idx, (cx, cy, shirt, accent) in enumerate(people):
        skin = (219, 174, 128, 255) if idx % 2 else (184, 132, 88, 255)
        hair = (43, 31, 24, 255)
        draw.ellipse([cx - 42, cy - 250, cx + 42, cy - 166], fill=skin)
        draw.pieslice([cx - 48, cy - 258, cx + 48, cy - 174], 180, 360, fill=hair)
        draw.rounded_rectangle([cx - 76, cy - 162, cx + 76, cy + 72], radius=38, fill=shirt)
        draw.rectangle([cx - 76, cy - 78, cx + 76, cy - 48], fill=accent)
        draw.line([cx - 65, cy - 130, cx - 158, cy - 34], fill=(37, 40, 37, 255), width=28)
        draw.line([cx + 65, cy - 130, cx + 158, cy - 34], fill=(37, 40, 37, 255), width=28)
        draw.ellipse([cx - 190, cy - 72, cx - 128, cy - 10], fill=(183, 255, 0, 255))
        draw.ellipse([cx + 128, cy - 72, cx + 190, cy - 10], fill=(18, 18, 18, 255))
        draw.line([cx - 44, cy + 70, cx - 75, cy + 245], fill=(45, 48, 45, 255), width=35)
        draw.line([cx + 44, cy + 70, cx + 75, cy + 245], fill=(45, 48, 45, 255), width=35)

    # camera/light texture, no text
    for x in range(-200, width, 260):
        draw.line([x, height, x + 520, 0], fill=(183, 255, 0, 35), width=18)
    img = img.filter(ImageFilter.SMOOTH_MORE)
    img.save(path, quality=92)


def make_logo(path):
    size = 900
    img = Image.new("RGBA", (size, size), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    cx = cy = size // 2
    points = []
    for i in range(10):
        radius = 330 if i % 2 == 0 else 155
        angle = -pi / 2 + i * pi / 5
        points.append((cx + cos(angle) * radius, cy + sin(angle) * radius))
    draw.polygon(points, fill=(183, 255, 0, 255))
    draw.ellipse([210, 210, 690, 690], fill=(18, 18, 18, 255))
    draw.polygon([(450, 260), (575, 610), (450, 535), (325, 610)], fill=(255, 255, 255, 255))
    img.save(path)


def text_box(slide, x, y, w, h, text, font, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rich_age_box(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    for value, color in [
        ("DAI ", BLACK),
        ("{{AGE_FROM}}", GREEN_DARK),
        (" AI ", BLACK),
        ("{{AGE_TO}}", GREEN_DARK),
        (" ANNI", BLACK),
    ]:
        r = p.add_run()
        r.text = value
        r.font.name = "Montserrat"
        r.font.size = Pt(21)
        r.font.bold = True
        r.font.color.rgb = color
    return box


def named_placeholder(slide, name, x, y, w, h, fill=LIGHT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(215, 221, 211)
    shape.line.width = Pt(1)
    return shape


def icon_line(slide, x, y, kind):
    group = []
    if kind == "pin":
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.14), Inches(0.14))
        s.fill.solid(); s.fill.fore_color.rgb = GREEN; s.line.color.rgb = GREEN
        group.append(s)
    elif kind == "phone":
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.16), Inches(0.24))
        s.fill.background(); s.line.color.rgb = GREEN; s.line.width = Pt(1.4)
        group.append(s)
    else:
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.18), Inches(0.18))
        s.fill.background(); s.line.color.rgb = GREEN; s.line.width = Pt(1.4)
        group.append(s)
    return group


def benefit_icon(slide, cx, y, kind):
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.18), Inches(y), Inches(0.36), Inches(0.36))
    circle.fill.background()
    circle.line.color.rgb = GREEN
    circle.line.width = Pt(1.4)
    if kind in {"DISCIPLINA", "FOCUS"}:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.04), Inches(y + 0.14), Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.color.rgb = GREEN
    elif kind == "SICUREZZA":
        shield = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, Inches(cx - 0.1), Inches(y + 0.09), Inches(0.2), Inches(0.2))
        shield.fill.background(); shield.line.color.rgb = GREEN; shield.line.width = Pt(1)
    elif kind == "FORZA":
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(cx - 0.12), Inches(y + 0.16), Inches(0.24), Inches(0.04))
        bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.color.rgb = GREEN
    elif kind == "RESISTENZA":
        bolt = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT, Inches(cx - 0.08), Inches(y + 0.07), Inches(0.16), Inches(0.22))
        bolt.fill.solid(); bolt.fill.fore_color.rgb = GREEN; bolt.line.color.rgb = GREEN
    else:
        star = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, Inches(cx - 0.09), Inches(y + 0.08), Inches(0.18), Inches(0.18))
        star.fill.background(); star.line.color.rgb = GREEN; star.line.width = Pt(1)


def create_template():
    prs = Presentation()
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(3.95), Inches(8.27), Inches(1.6)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = LIGHT
    slide.shapes[-1].line.fill.background()

    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(8.8), Inches(8.27), Inches(2.89)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = RGBColor(248, 249, 246)
    slide.shapes[-1].line.fill.background()

    # Header / Hero
    logo_box = text_box(slide, 0.45, 0.36, 2.4, 0.32, "VIRTUS COMBAT ACADEMY", "Montserrat", 8.5, BLACK, True)
    logo_box.name = "BRAND_TEXT"
    mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(0.72), Inches(0.86), Inches(0.08))
    mark.fill.solid(); mark.fill.fore_color.rgb = GREEN; mark.line.fill.background()

    title1 = text_box(slide, 0.43, 0.98, 3.05, 0.92, "{{TITLE}}", "Anton", 58, GREEN_DARK, True)
    title1.name = "TEXT_TITLE"
    title2 = text_box(slide, 0.43, 1.72, 3.45, 1.0, "{{SUBTITLE}}", "Anton", 52, BLACK, True)
    title2.name = "TEXT_SUBTITLE"
    rich_age_box(slide, 0.47, 2.72, 3.25, 0.34).name = "TEXT_AGE"
    claim = text_box(slide, 0.48, 3.14, 2.85, 0.42, "{{CLAIM}}", "Montserrat", 10.5, ANTHRACITE, True)
    claim.name = "TEXT_CLAIM"

    hero = named_placeholder(slide, "IMAGE_HERO", 3.55, 0.28, 4.28, 3.42)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(3.28), Inches(0.28), Inches(0.34), Inches(3.42))
    strip.fill.solid(); strip.fill.fore_color.rgb = GREEN; strip.line.fill.background()

    # Gallery
    gallery_y = 4.15
    named_placeholder(slide, "IMAGE_GALLERY_1", 0.48, gallery_y, 2.22, 1.13)
    named_placeholder(slide, "IMAGE_GALLERY_2", 3.02, gallery_y, 2.22, 1.13)
    named_placeholder(slide, "IMAGE_GALLERY_3", 5.56, gallery_y, 2.22, 1.13)
    for x in [2.76, 5.3]:
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(x), Inches(gallery_y - 0.02), Inches(0.22), Inches(1.17))
        shp.fill.solid(); shp.fill.fore_color.rgb = GREEN; shp.line.fill.background()

    # Benefits
    text_box(slide, 0.75, 5.82, 6.78, 0.42, "COSA TI PORTERAI DENTRO", "Anton", 27, BLACK, True, PP_ALIGN.CENTER)
    benefits = [
        ("DISCIPLINA", "{{BENEFIT_1_TEXT}}"),
        ("SICUREZZA", "{{BENEFIT_2_TEXT}}"),
        ("FORZA", "{{BENEFIT_3_TEXT}}"),
        ("RESISTENZA", "{{BENEFIT_4_TEXT}}"),
        ("FOCUS", "{{BENEFIT_5_TEXT}}"),
        ("AUTOSTIMA", "{{BENEFIT_6_TEXT}}"),
    ]
    start_x = 0.34
    col_w = 1.27
    for i, (title, desc) in enumerate(benefits):
        x = start_x + i * col_w
        cx = x + col_w / 2
        benefit_icon(slide, cx, 6.46, title)
        text_box(slide, x, 6.91, col_w, 0.2, "{{BENEFIT_%d_TITLE}}" % (i + 1), "Montserrat", 7.8, BLACK, True, PP_ALIGN.CENTER)
        text_box(slide, x + 0.06, 7.19, col_w - 0.12, 0.82, desc, "Montserrat", 6.7, MUTED, False, PP_ALIGN.CENTER)

    # CTA
    brush = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(0.43), Inches(8.9), Inches(3.62), Inches(0.72))
    brush.fill.solid(); brush.fill.fore_color.rgb = BLACK; brush.line.fill.background()
    text_box(slide, 0.7, 9.0, 3.18, 0.42, "{{CTA}}", "Anton", 29, WHITE, True, PP_ALIGN.CENTER)
    text_box(slide, 0.53, 9.77, 3.25, 0.44, "{{CTA_SUBTITLE}}", "Montserrat", 12.5, BLACK, True)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(4.18), Inches(9.04), Inches(1.38), Inches(0.66))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = GREEN; arrow.line.fill.background()

    # Footer
    footer_x = 5.45
    icon_line(slide, footer_x, 9.12, "pin")
    text_box(slide, footer_x + 0.25, 9.08, 1.8, 0.2, "{{ADDRESS}}", "Montserrat", 8.5, BLACK, True)
    icon_line(slide, footer_x, 9.52, "phone")
    text_box(slide, footer_x + 0.25, 9.5, 1.8, 0.2, "{{PHONE}}", "Montserrat", 8.5, BLACK, True)
    icon_line(slide, footer_x, 9.92, "instagram")
    text_box(slide, footer_x + 0.25, 9.9, 1.8, 0.2, "{{INSTAGRAM}}", "Montserrat", 8.5, BLACK, True)
    named_placeholder(slide, "LOGO", 7.0, 9.16, 0.72, 0.72, WHITE)
    text_box(slide, 6.24, 10.02, 1.55, 0.32, "VIRTUS\nCOMBAT ACADEMY", "Montserrat", 6.6, BLACK, True, PP_ALIGN.RIGHT)

    template_path = TEMPLATE_DIR / "lotta_pugilato_light_vertical.pptx"
    prs.save(template_path)
    return template_path


def create_json():
    data = """{
  "template": "lotta_pugilato_light_vertical.pptx",
  "output_name": "lotta_pugilato_10_14",
  "texts": {
    "TITLE": "LOTTA",
    "SUBTITLE": "PUGILATO",
    "AGE_FROM": "10",
    "AGE_TO": "14",
    "CLAIM": "ENERGIA, AMICIZIA, DISCIPLINA E SICUREZZA.",
    "CTA": "PROVA GRATUITA",
    "CTA_SUBTITLE": "LA PRIMA LEZIONE TI ASPETTIAMO!",
    "PHONE": "351 8899843",
    "INSTAGRAM": "@virtus_group_",
    "ADDRESS": "Via Corfu 71",
    "BENEFIT_1_TITLE": "DISCIPLINA",
    "BENEFIT_1_TEXT": "Impari il valore dell'impegno e della costanza.",
    "BENEFIT_2_TITLE": "SICUREZZA",
    "BENEFIT_2_TEXT": "Aumenti la sicurezza in te stesso.",
    "BENEFIT_3_TITLE": "FORZA",
    "BENEFIT_3_TEXT": "Sviluppi forza fisica e controllo.",
    "BENEFIT_4_TITLE": "RESISTENZA",
    "BENEFIT_4_TEXT": "Migliori fiato ed energia.",
    "BENEFIT_5_TITLE": "FOCUS",
    "BENEFIT_5_TEXT": "Alleni concentrazione e mentalita.",
    "BENEFIT_6_TITLE": "AUTOSTIMA",
    "BENEFIT_6_TEXT": "Cresci e scopri il meglio di te."
  },
  "images": {
    "IMAGE_HERO": "assets/images/lotta_pugilato_hero.jpg",
    "IMAGE_GALLERY_1": "assets/images/lotta_pugilato_drill.jpg",
    "IMAGE_GALLERY_2": "assets/images/lotta_pugilato_lotta.jpg",
    "IMAGE_GALLERY_3": "assets/images/lotta_pugilato_sacco.jpg",
    "LOGO": "assets/images/logo_virtus_combat_academy.png"
  }
}
"""
    path = DATA_DIR / "lotta_pugilato_10_14.json"
    path.write_text(data, encoding="utf-8")
    return path


def main():
    ensure_dirs()
    make_photo(ASSET_DIR / "lotta_pugilato_hero.jpg", "hero")
    make_photo(ASSET_DIR / "lotta_pugilato_drill.jpg", "drill")
    make_photo(ASSET_DIR / "lotta_pugilato_lotta.jpg", "wrestling")
    make_photo(ASSET_DIR / "lotta_pugilato_sacco.jpg", "bag")
    make_logo(ASSET_DIR / "logo_virtus_combat_academy.png")
    template_path = create_template()
    json_path = create_json()
    print(template_path)
    print(json_path)


if __name__ == "__main__":
    main()
